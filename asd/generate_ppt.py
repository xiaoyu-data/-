from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# ============================================================
#  配色方案
# ============================================================
COLOR_BG_TOP      = RGBColor(235, 245, 255)
COLOR_BG_BOTTOM   = RGBColor(210, 230, 250)
COLOR_CARD_BG     = RGBColor(255, 255, 255)
COLOR_ICON_CIRCLE = RGBColor(0, 120, 220)
COLOR_ICON_INNER  = RGBColor(255, 255, 255)
COLOR_TITLE       = RGBColor(30, 40, 60)
COLOR_DESC        = RGBColor(120, 130, 150)
COLOR_DIVIDER     = RGBColor(200, 210, 220)

# ============================================================
#  幻灯片尺寸 (16:9)
# ============================================================
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ============================================================
#  功能数据
# ============================================================
features = [
    ("全面可视", "资源状态一目了然",   "eye"),
    ("智能诊断", "故障快速定位",       "pulse"),
    ("高效运维", "闭环流程，降本增效", "gear"),
    ("安全可靠", "多维防护，稳定运行", "shield"),
]

# ============================================================
#  布局参数
# ============================================================
CARD_WIDTH   = Inches(2.6)
CARD_HEIGHT  = Inches(2.2)
ICON_SIZE    = Inches(0.9)
CARD_GAP     = Inches(0.15)
DIVIDER_H    = Inches(1.4)

# 计算居中起始位置
total_width = len(features) * CARD_WIDTH + (len(features) - 1) * CARD_GAP
LEFT_START  = (SLIDE_WIDTH.inches - total_width) / 2
LEFT_START  = Inches(LEFT_START)
TOP         = Inches(2.5)


# ============================================================
#  辅助函数
# ============================================================
def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_text_style(paragraph, text, font_size, font_color, bold=False, font_name="Microsoft YaHei"):
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_shape_fill(shape, fill_color)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_background_gradient(slide, color_top, color_bottom):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 270
    fill.gradient_stops[0].color.rgb = color_top
    fill.gradient_stops[1].color.rgb = color_bottom


def _to_inches(val):
    if isinstance(val, (int, float)):
        return Inches(val)
    return val


# ============================================================
#  图标绘制函数
# ============================================================
def draw_eye_icon(slide, cx, cy, size):
    cx = _to_inches(cx)
    cy = _to_inches(cy)
    size = _to_inches(size)
    eye = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - size/2, cy - size/2.5, size, size/1.25
    )
    set_shape_fill(eye, COLOR_ICON_INNER)
    eye.line.color.rgb = COLOR_ICON_INNER
    eye.line.width = Pt(2)
    pupil = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - size/6, cy - size/6, size/3, size/3
    )
    set_shape_fill(pupil, COLOR_ICON_CIRCLE)
    pupil.line.fill.background()


def draw_pulse_icon(slide, cx, cy, size):
    cx = _to_inches(cx)
    cy = _to_inches(cy)
    size = _to_inches(size)
    s = size.inches
    cx_f = cx.inches
    cy_f = cy.inches
    from pptx.enum.shapes import MSO_CONNECTOR
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - size/2, cy - size/2.5, size, size/1.25
    )
    set_shape_fill(shape, COLOR_ICON_INNER)
    shape.line.fill.background()
    pts = [
        (cx_f - s*0.35, cy_f),
        (cx_f - s*0.2,  cy_f),
        (cx_f - s*0.1,  cy_f - s*0.2),
        (cx_f,          cy_f + s*0.15),
        (cx_f + s*0.1,  cy_f - s*0.25),
        (cx_f + s*0.2,  cy_f),
        (cx_f + s*0.35, cy_f),
    ]
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = COLOR_ICON_CIRCLE
        line.line.width = Pt(2.5)


def draw_gear_icon(slide, cx, cy, size):
    cx = _to_inches(cx)
    cy = _to_inches(cy)
    size = _to_inches(size)
    gear = slide.shapes.add_shape(
        MSO_SHAPE.GEAR_6, cx - size/2, cy - size/2, size, size
    )
    set_shape_fill(gear, COLOR_ICON_INNER)
    gear.line.color.rgb = COLOR_ICON_INNER
    gear.line.width = Pt(2)
    center = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - size/6, cy - size/6, size/3, size/3
    )
    set_shape_fill(center, COLOR_ICON_CIRCLE)
    center.line.fill.background()


def draw_shield_icon(slide, cx, cy, size):
    cx = _to_inches(cx)
    cy = _to_inches(cy)
    size = _to_inches(size)
    shield = slide.shapes.add_shape(
        MSO_SHAPE.PENTAGON, cx - size/2, cy - size/2, size, size
    )
    set_shape_fill(shield, COLOR_ICON_INNER)
    shield.line.color.rgb = COLOR_ICON_INNER
    shield.line.width = Pt(2)
    check = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, cx - size/3, cy - size/4, size/1.5, size/2
    )
    set_shape_fill(check, COLOR_ICON_CIRCLE)
    check.line.fill.background()
    check.rotation = 90


ICON_DRAWERS = {
    "eye":    draw_eye_icon,
    "pulse":  draw_pulse_icon,
    "gear":   draw_gear_icon,
    "shield": draw_shield_icon,
}


# ============================================================
#  主程序
# ============================================================
prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

slide = prs.slides.add_slide(prs.slide_layouts[6])

# 添加渐变背景
add_background_gradient(slide, COLOR_BG_TOP, COLOR_BG_BOTTOM)

# 遍历功能数据，绘制卡片
for i, (title, desc, icon_type) in enumerate(features):
    left = LEFT_START + i * (CARD_WIDTH + CARD_GAP)
    card_center_x = left + CARD_WIDTH / 2

    # ---------- 卡片背景（圆角矩形）----------
    card = add_rounded_rect(
        slide, left, TOP, CARD_WIDTH, CARD_HEIGHT,
        fill_color=COLOR_CARD_BG
    )

    # ---------- 图标圆形背景 ----------
    icon_bg = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        card_center_x - ICON_SIZE / 2,
        TOP + Inches(0.25),
        ICON_SIZE, ICON_SIZE
    )
    set_shape_fill(icon_bg, COLOR_ICON_CIRCLE)
    icon_bg.line.fill.background()

    # ---------- 绘制图标 ----------
    icon_cx = card_center_x
    icon_cy = TOP + Inches(0.25) + ICON_SIZE / 2
    draw_fn = ICON_DRAWERS.get(icon_type)
    if draw_fn:
        draw_fn(slide, icon_cx, icon_cy, ICON_SIZE * 0.5)

    # ---------- 标题 ----------
    title_box = slide.shapes.add_textbox(
        left, TOP + Inches(1.15), CARD_WIDTH, Inches(0.4)
    )
    tf_title = title_box.text_frame
    tf_title.word_wrap = False
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    set_text_style(p_title, title, font_size=16, font_color=COLOR_TITLE, bold=True)

    # ---------- 描述 ----------
    desc_box = slide.shapes.add_textbox(
        left, TOP + Inches(1.55), CARD_WIDTH, Inches(0.4)
    )
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.alignment = PP_ALIGN.CENTER
    set_text_style(p_desc, desc, font_size=12, font_color=COLOR_DESC)

    # ---------- 分隔线（除最后一个）----------
    if i < len(features) - 1:
        div_left = left + CARD_WIDTH + CARD_GAP / 2 - Inches(0.015)
        div_top  = TOP + (CARD_HEIGHT - DIVIDER_H) / 2
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            div_left, div_top,
            Inches(0.03), DIVIDER_H
        )
        set_shape_fill(divider, COLOR_DIVIDER)
        divider.line.fill.background()

# ============================================================
#  保存
# ============================================================
output_path = r"C:\Users\熊凯鹏\Desktop\功能点展示.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
