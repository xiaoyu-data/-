from pptx import Presentation

prs = Presentation(r"C:\Users\熊凯鹏\Desktop\功能点展示.pptx")
slide = prs.slides[0]

print(f"幻灯片尺寸: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print(f"形状数量: {len(slide.shapes)}")
print()
print("形状列表:")
for s in slide.shapes:
    name = s.name
    stype = str(s.shape_type).replace("MSO_SHAPE_TYPE.", "")
    pos = f"({s.left.inches:.2f}, {s.top.inches:.2f})"
    size = f"({s.width.inches:.2f}, {s.height.inches:.2f})"
    print(f"  - {stype:12s} | {name:20s} | pos={pos:18s} | size={size}")

print()
print("文本内容:")
for s in slide.shapes:
    if s.has_text_frame:
        text = s.text_frame.text.strip()
        if text:
            print(f"  - {text}")
