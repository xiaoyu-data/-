from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# 测试1: 纯色背景
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(210, 230, 250)

# 测试2: 简单矩形
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape.line.fill.background()

# 测试3: 文本
tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = "Hello World"
p.alignment = PP_ALIGN.LEFT
run = p.runs[0]
run.font.size = Pt(24)
run.font.color.rgb = RGBColor(0, 0, 0)

prs.save(r"C:\Users\熊凯鹏\Desktop\test_minimal.pptx")
print("test_minimal.pptx 已保存")
